"""Generate PowerPoint slide deck for Senior Staff Engineer assessment presentation."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).parent / "Separation_of_Concerns_Assessment.pptx"

# Slide 1: Title
SLIDE1_TITLE = "Code Review: Separation of Concerns"
SLIDE1_SUBTITLE = "S504_TriviaMaze_TeamDelta — Senior Staff Engineer Assessment"
SLIDE1_NOTES = (
    "Today I'll present a brief code review of our Trivia Maze codebase, focusing on "
    "separation of concerns between the maze domain, persistence layer, and presentation. "
    "I'll cover what we did well, one violation we found, and a concrete fix."
)

# Slide 2: Assessment
SLIDE2_TITLE = "Separation of Concerns Assessment"
SLIDE2_BULLETS = [
    "Maze and DB are fully decoupled — neither imports from the other",
    "main.py is the sole translation layer (Position ↔ primitives)",
    "I/O violation: main.py line 406 calls input() directly for replay prompt",
    "Per interfaces.md: all print() and input() must reside in view.py only",
    "View imports only from maze; GUI modules correctly use QuizMazeGame",
]
SLIDE2_NOTES = (
    "The codebase maintains strict boundaries. maze.py has zero project imports—only stdlib. "
    "db.py uses only primitives like player_row, player_col, visited_cells as lists—no maze types. "
    "main.py handles all translation. However, we found one violation: when the player wins, "
    "main.py asks 'Play again?' using a raw input() call instead of delegating to view. "
    "Our interfaces document says view.py owns all I/O. We need to fix that."
)

# Slide 3: Refactoring
SLIDE3_TITLE = "Concrete Refactoring Recommendation"
SLIDE3_BULLETS = [
    "Replace input('\\nPlay again? (y/n): ') in main.py with view.prompt_replay()",
    "view.prompt_replay() already exists with the themed prompt",
    "Restores the I/O rule; centralizes all prompts in view",
    "One-line change: if view.prompt_replay(): instead of manual input/parse",
]
SLIDE3_NOTES = (
    "The fix is simple. view.py already has prompt_replay() which asks 'Venture into the forest again?' "
    "and returns True for yes. We just need to replace the raw input block in main.py with a call "
    "to view.prompt_replay(). That restores our I/O constraint, keeps main focused on orchestration, "
    "and makes the codebase consistent with how we handle prompt_name and prompt_resume."
)

# Slide 4: Summary
SLIDE4_TITLE = "Summary"
SLIDE4_BULLETS = [
    "Architecture is sound — maze and db stay decoupled",
    "One I/O leak: replay prompt in main.py",
    "Fix: use view.prompt_replay() — one line",
]
SLIDE4_NOTES = (
    "To wrap up: our separation of concerns is strong. The only issue is the replay prompt. "
    "The fix is a one-line change. Questions?"
)


def add_title_slide(prs: Presentation) -> None:
    layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = SLIDE1_TITLE
    slide.placeholders[1].text = SLIDE1_SUBTITLE
    slide.notes_slide.notes_text_frame.text = SLIDE1_NOTES


def add_content_slide(prs: Presentation, title: str, bullets: list[str], notes: str) -> None:
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
    slide.notes_slide.notes_text_frame.text = notes


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_content_slide(prs, SLIDE2_TITLE, SLIDE2_BULLETS, SLIDE2_NOTES)
    add_content_slide(prs, SLIDE3_TITLE, SLIDE3_BULLETS, SLIDE3_NOTES)
    add_content_slide(prs, SLIDE4_TITLE, SLIDE4_BULLETS, SLIDE4_NOTES)

    prs.save(OUTPUT)
    print(f"Created: {OUTPUT}")


if __name__ == "__main__":
    main()
